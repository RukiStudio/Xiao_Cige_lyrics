# -*- coding: utf-8 -*-
"""小词格 AI 歌词生成器 - 介绍 PPT 生成脚本"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============ 全局配置 ============

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# 配色方案
C_BG_DARK = RGBColor(0x0F, 0x17, 0x2A)       # 深蓝背景
C_BG_MID = RGBColor(0x1E, 0x29, 0x3B)        # 次级背景
C_PRIMARY = RGBColor(0x81, 0x8C, 0xF8)        # 紫色主色
C_ACCENT = RGBColor(0x63, 0x66, 0xF1)        # 深紫
C_TEXT = RGBColor(0xFF, 0xFF, 0xFF)           # 白色文字
C_TEXT_DIM = RGBColor(0x94, 0xA3, 0xB8)       # 次级文字
C_TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)    # 更暗文字
C_GREEN = RGBColor(0x34, 0xD3, 0x99)         # 绿色强调
C_ORANGE = RGBColor(0xFB, 0x92, 0x3C)        # 橙色强调
C_BORDER = RGBColor(0x33, 0x41, 0x55)        # 边框色
C_CARD = RGBColor(0x1E, 0x29, 0x3B)          # 卡片背景

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT_PATH = os.path.join(ROOT, "docs", "小词格介绍.pptx")
SHOT_DIR = os.path.join(ROOT, "docs", "screenshots")

# ============ 工具函数 ============

def set_slide_bg(slide, color):
    """设置幻灯片纯色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=C_TEXT, alignment=PP_ALIGN.LEFT,
                font_name="Microsoft YaHei"):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=C_TEXT, bullet_color=C_PRIMARY):
    """添加项目符号列表"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # 项目符号用彩色方框
        run1 = p.add_run()
        run1.text = "● "
        run1.font.size = Pt(font_size)
        run1.font.color.rgb = bullet_color
        run1.font.bold = True
        run1.font.name = "Microsoft YaHei"
        # 文本
        run2 = p.add_run()
        run2.text = item
        run2.font.size = Pt(font_size)
        run2.font.color.rgb = color
        run2.font.name = "Microsoft YaHei"
        p.space_after = Pt(8)
    return txBox


def add_card(slide, left, top, width, height, fill_color=C_CARD,
             border_color=None):
    """添加圆角卡片"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_page_header(slide, title, subtitle=None, accent_color=C_PRIMARY):
    """统一页面头部"""
    # 左侧色条
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.6),
        Inches(0.08), Inches(0.6)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()

    # 标题
    add_textbox(slide, Inches(0.85), Inches(0.55),
                Inches(10), Inches(0.7), title,
                font_size=32, bold=True, color=C_TEXT)
    if subtitle:
        add_textbox(slide, Inches(0.85), Inches(1.25),
                    Inches(10), Inches(0.5), subtitle,
                    font_size=16, color=C_TEXT_DIM)

    # 底部装饰线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(7.0),
        Inches(12.1), Emu(30000)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = C_BORDER
    line.line.fill.background()


def add_page_number(slide, num, total):
    """添加页码"""
    add_textbox(slide, Inches(12.2), Inches(7.0),
                Inches(1), Inches(0.4),
                f"{num} / {total}",
                font_size=11, color=C_TEXT_MUTED,
                alignment=PP_ALIGN.RIGHT)


def add_image_safe(slide, img_path, left, top, width, height):
    """安全添加图片（文件不存在则显示占位文字）"""
    if os.path.exists(img_path):
        try:
            return slide.shapes.add_picture(img_path, left, top, width, height)
        except Exception:
            pass
    # 占位
    box = add_card(slide, left, top, width, height,
                   RGBColor(0x1A, 0x22, 0x36), C_BORDER)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "📷 截图待补充"
    p.font.size = Pt(12)
    p.font.color.rgb = C_TEXT_MUTED
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(40)
    return box


TOTAL_SLIDES = 13

# ============ 开始构建 ============

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]

# ==================== Slide 1: 封面 ====================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)

# 装饰大色块
deco = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
    Inches(5), Inches(7.5)
)
deco.fill.solid()
deco.fill.fore_color.rgb = C_ACCENT
deco.line.fill.background()

# 右侧装饰圆
circle = slide.shapes.add_shape(
    MSO_SHAPE.OVAL, Inches(10), Inches(-2),
    Inches(6), Inches(6)
)
circle.fill.solid()
circle.fill.fore_color.rgb = C_ACCENT
circle.fill.transparency = 0.7
circle.line.fill.background()

circle2 = slide.shapes.add_shape(
    MSO_SHAPE.OVAL, Inches(11), Inches(4),
    Inches(3), Inches(3)
)
circle2.fill.solid()
circle2.fill.fore_color.rgb = C_PRIMARY
circle2.fill.transparency = 0.6
circle2.line.fill.background()

# 标题
add_textbox(slide, Inches(5.5), Inches(2.2),
            Inches(7), Inches(1),
            "小词格", font_size=64, bold=True, color=C_TEXT)

add_textbox(slide, Inches(5.5), Inches(3.2),
            Inches(7), Inches(0.7),
            "AI 歌词生成器", font_size=36, bold=True, color=C_PRIMARY)

add_textbox(slide, Inches(5.5), Inches(4.1),
            Inches(7), Inches(0.5),
            "纯大模型驱动 · 多轮自我迭代 · 对抗评审修改",
            font_size=18, color=C_TEXT_DIM)

# 底部信息
add_textbox(slide, Inches(5.5), Inches(6.3),
            Inches(7), Inches(0.4),
            "作者哔哩哔哩：space.bilibili.com/527121484",
            font_size=14, color=C_TEXT_MUTED)

add_textbox(slide, Inches(5.5), Inches(6.8),
            Inches(7), Inches(0.4),
            "2026  ·  Windows 桌面版 v1.0.0",
            font_size=12, color=C_TEXT_MUTED)

# ==================== Slide 2: 目录 ====================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)
add_page_header(slide, "目录", "CONTENTS")

items = [
    ("01", "项目概述", "是什么 · 解决什么问题"),
    ("02", "核心架构", "多智能体协作流程"),
    ("03", "词格系统", "0 / 1 / / / + 标记体系"),
    ("04", "英语歌词模式", "多音节词适配方案"),
    ("05", "评审严格度", "动态覆盖 Critic 提示词"),
    ("06", "技术实现", "Next.js · SSE · 重试机制"),
    ("07", "桌面安装包", "Electron 打包方案"),
    ("08", "演示与总结", "实机截图 · 未来规划"),
]

for i, (num, title, desc) in enumerate(items):
    row = i // 2
    col = i % 2
    left = Inches(0.8 + col * 6.2)
    top = Inches(1.7 + row * 1.3)

    card = add_card(slide, left, top, Inches(5.8), Inches(1.1), C_CARD, C_BORDER)

    # 数字
    add_textbox(slide, left + Inches(0.2), top + Inches(0.15),
                Inches(0.8), Inches(0.8), num,
                font_size=28, bold=True, color=C_PRIMARY)
    # 标题
    add_textbox(slide, left + Inches(1.0), top + Inches(0.15),
                Inches(4.5), Inches(0.5), title,
                font_size=20, bold=True, color=C_TEXT)
    # 描述
    add_textbox(slide, left + Inches(1.0), top + Inches(0.6),
                Inches(4.5), Inches(0.4), desc,
                font_size=13, color=C_TEXT_DIM)

add_page_number(slide, 2, TOTAL_SLIDES)

# ==================== Slide 3: 项目概述 ====================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)
add_page_header(slide, "项目概述", "PROJECT OVERVIEW")

# 左侧大图
img_path = os.path.join(SHOT_DIR, "01-main-ui.png")
if os.path.exists(img_path):
    add_image_safe(slide, img_path, Inches(0.6), Inches(1.8),
                   Inches(7.5), Inches(4.8))
else:
    add_card(slide, Inches(0.6), Inches(1.8), Inches(7.5), Inches(4.8),
             C_CARD, C_BORDER)
    add_textbox(slide, Inches(0.6), Inches(4.0), Inches(7.5), Inches(0.5),
                "📷 主界面截图", font_size=14, color=C_TEXT_MUTED,
                alignment=PP_ALIGN.CENTER)

# 右侧内容
right_left = Inches(8.4)
right_w = Inches(4.5)

add_textbox(slide, right_left, Inches(1.9), right_w, Inches(0.6),
            "一句话介绍", font_size=22, bold=True, color=C_PRIMARY)

add_textbox(slide, right_left, Inches(2.5), right_w, Inches(1.2),
            "基于多智能体对抗式迭代的 AI 歌词生成器，支持中文/英文，严格遵循词格约束。",
            font_size=15, color=C_TEXT, font_name="Microsoft YaHei")

# 核心特性
add_textbox(slide, right_left, Inches(3.7), right_w, Inches(0.4),
            "核心特性", font_size=18, bold=True, color=C_PRIMARY)

features = [
    "四个 AI 智能体协作生成",
    "逐字匹配词格，支持延音 (+)",
    "英语歌词按音节数校验",
    "评审严格度 1-10 级可调",
    "SSE 流式实时进度",
    "Windows 桌面安装包",
]
add_bullet_list(slide, right_left, Inches(4.2), right_w, Inches(3),
                features, font_size=13)

add_page_number(slide, 3, TOTAL_SLIDES)

# ==================== Slide 4: 多智能体架构 ====================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)
add_page_header(slide, "核心架构", "多智能体对抗式迭代")

# 四个智能体卡片
agents = [
    ("Generator", "生成器", "根据提示词与词格\n生成初始歌词", C_PRIMARY),
    ("Critic", "评审员", "七维度打分\n标记问题", C_ORANGE),
    ("Reviser", "修改器", "按评审意见修改\n严格保持字数", C_GREEN),
    ("Interpreter", "词解师", "生成词解与\n风格关键词", RGBColor(0xF4, 0x72, 0xB6)),
]

for i, (name, cn, desc, color) in enumerate(agents):
    left = Inches(0.6 + i * 3.15)
    top = Inches(2.0)
    w = Inches(2.95)
    h = Inches(2.5)

    card = add_card(slide, left, top, w, h, C_CARD, color)

    # 色带
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top,
        w, Emu(80000)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_textbox(slide, left + Inches(0.2), top + Inches(0.25),
                w - Inches(0.4), Inches(0.5),
                name, font_size=20, bold=True, color=color)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.75),
                w - Inches(0.4), Inches(0.4),
                cn, font_size=14, color=C_TEXT_DIM)
    add_textbox(slide, left + Inches(0.2), top + Inches(1.2),
                w - Inches(0.4), Inches(1.2),
                desc, font_size=13, color=C_TEXT)

# 流程箭头
for i in range(3):
    arrow_left = Inches(3.55 + i * 3.15)
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, arrow_left, Inches(3.05),
        Inches(0.25), Inches(0.3)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = C_BORDER
    arrow.line.fill.background()

# 流程说明
add_textbox(slide, Inches(0.6), Inches(4.8), Inches(12), Inches(0.5),
            "生成 → 评审 → 修改 → 达标？ → 词解输出",
            font_size=20, bold=True, color=C_PRIMARY, alignment=PP_ALIGN.CENTER)

# 循环说明
loop_items = [
    "内容迭代：最多 3 轮（生成→评审→修改 主循环）",
    "字数修复：纯字数问题额外最多 5 轮，不计入主迭代",
    "自动重试：LLM 调用超时自动重试最多 6 次",
    "用户可随时：点击「停止生成」中止流程",
]
add_bullet_list(slide, Inches(0.6), Inches(5.3), Inches(12), Inches(2),
                loop_items, font_size=14)

add_page_number(slide, 4, TOTAL_SLIDES)

# ==================== Slide 5: 词格系统 ====================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)
add_page_header(slide, "词格系统", "高级编辑模式 · 0 / 1 / / / + 标记")

# 左侧表格
table_data = [
    ["标记", "含义", "计入字数", "说明"],
    ["0 或汉字", "普通音节", "✅", "一个汉字 = 一个音节"],
    ["1", "重音音节", "✅", "需强调的音节"],
    ["/", "停顿", "❌", "气口 / 断句"],
    ["+", "延音", "❌", "拉长前一音节"],
]

rows = len(table_data)
cols = 4
table = slide.shapes.add_table(
    rows, cols, Inches(0.6), Inches(1.8),
    Inches(8), Inches(2.5)
).table

col_widths = [Inches(1.5), Inches(1.8), Inches(1.5), Inches(3.2)]
for i, w in enumerate(col_widths):
    table.columns[i].width = w

for r in range(rows):
    for c in range(cols):
        cell = table.cell(r, c)
        cell.text = table_data[r][c]
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(14)
                run.font.name = "Microsoft YaHei"
                if r == 0:
                    run.font.bold = True
                    run.font.color.rgb = C_TEXT
                else:
                    run.font.color.rgb = C_TEXT_DIM if r % 2 == 0 else C_TEXT
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_PRIMARY if r == 0 else (C_CARD if r % 2 == 0 else RGBColor(0x25, 0x33, 0x49))

# 右侧截图
img_path = os.path.join(SHOT_DIR, "02-advanced-editor.png")
add_textbox(slide, Inches(8.8), Inches(1.8), Inches(4), Inches(0.4),
            "实机截图", font_size=16, bold=True, color=C_PRIMARY)
add_image_safe(slide, img_path, Inches(8.8), Inches(2.3),
               Inches(4), Inches(3))

# 底部说明
img_path2 = os.path.join(SHOT_DIR, "03-rhythm-tags.png")
add_textbox(slide, Inches(0.6), Inches(4.5), Inches(6), Inches(0.4),
            "+ 延音标记示例", font_size=16, bold=True, color=C_PRIMARY)
add_image_safe(slide, img_path2, Inches(0.6), Inches(5.0),
               Inches(6), Inches(2))

# 延音说明
add_textbox(slide, Inches(6.8), Inches(4.5), Inches(6), Inches(0.4),
            "延音在歌词中的体现", font_size=16, bold=True, color=C_PRIMARY)
hold_notes = [
    "中文：拉长字尾 呀 / 哦 / 啊～",
    "英文：拉长元音 yeaaah / hold ooon",
    "⚠️ 不得为延音新增单词",
]
add_bullet_list(slide, Inches(6.8), Inches(5.0), Inches(6), Inches(2),
                hold_notes, font_size=14)

add_page_number(slide, 5, TOTAL_SLIDES)

# ==================== Slide 6: 英语歌词模式 ====================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)
add_page_header(slide, "英语歌词模式", "多音节词适配 · syllable-based")

# 左侧
add_textbox(slide, Inches(0.6), Inches(1.7), Inches(6), Inches(0.5),
            "模式说明", font_size=20, bold=True, color=C_PRIMARY)

add_textbox(slide, Inches(0.6), Inches(2.3), Inches(6), Inches(1),
            "歌词用英文创作，评审 / 词解 / 修改说明仍为中文。",
            font_size=15, color=C_TEXT)

add_textbox(slide, Inches(0.6), Inches(3.3), Inches(6), Inches(0.4),
            "核心规则", font_size=18, bold=True, color=C_PRIMARY)

rules = [
    "按英语音节数（syllable）匹配词格，非单词数",
    "beautiful = 3 音节（beau-ti-ful）",
    "lonely = 2 音节（one-ly）",
    "love = 1 音节",
    "押韵在英文元音层面判定（rain / pain ✓）",
    "延音拉长前词末尾元音，不得新增单词",
]
add_bullet_list(slide, Inches(0.6), Inches(3.8), Inches(6), Inches(3),
                rules, font_size=13)

# 右侧截图
img_path = os.path.join(SHOT_DIR, "04-english-mode.png")
add_textbox(slide, Inches(7), Inches(1.7), Inches(5.5), Inches(0.4),
            "实机截图", font_size=16, bold=True, color=C_PRIMARY)
add_image_safe(slide, img_path, Inches(7), Inches(2.2),
               Inches(5.8), Inches(4.5))

add_page_number(slide, 6, TOTAL_SLIDES)

# ==================== Slide 7: 评审严格度 ====================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)
add_page_header(slide, "评审严格度", "1-10 级动态覆盖 Critic 提示词")

# 四个模式对比
modes = [
    ("1-3 宽松", "鼓励为主", "单句不符 ≤ 8 分\n平庸 ≤ 8 分\n套话 ≤ 6 分",
     "达标即可", C_GREEN),
    ("4-6 标准", "客观评审", "单句不符 ≤ 6 分\n平庸 ≤ 7 分\n套话 ≤ 5 分",
     "按阈值", C_PRIMARY),
    ("7-8 严格", "挑剔制作人", "单句不符 ≤ 4 分\n平庸 ≤ 6 分\n套话 ≤ 4 分",
     "有 high issue 不通过", C_ORANGE),
    ("9-10 极严", "顶级评审", "一字不符 → 0 分\n平庸 ≤ 4 分\n套话 ≤ 2 分",
     "structure=10 且无 high", RGBColor(0xEF, 0x44, 0x44)),
]

for i, (label, sub, rules_text, pass_text, color) in enumerate(modes):
    left = Inches(0.6 + i * 3.15)
    top = Inches(1.9)
    w = Inches(2.95)
    h = Inches(3.5)

    card = add_card(slide, left, top, w, h, C_CARD, color)

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, w, Emu(80000)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_textbox(slide, left + Inches(0.2), top + Inches(0.25),
                w - Inches(0.4), Inches(0.5),
                label, font_size=18, bold=True, color=color)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.7),
                w - Inches(0.4), Inches(0.4),
                sub, font_size=13, color=C_TEXT_DIM)
    add_textbox(slide, left + Inches(0.2), top + Inches(1.15),
                w - Inches(0.4), Inches(1.5),
                rules_text, font_size=12, color=C_TEXT)
    add_textbox(slide, left + Inches(0.2), top + Inches(2.7),
                w - Inches(0.4), Inches(0.4),
                "passed: " + pass_text, font_size=11,
                bold=True, color=color)

# 底部说明
add_textbox(slide, Inches(0.6), Inches(5.7), Inches(12), Inches(0.4),
            "严格度不只是附加描述，而是直接覆盖 Critic 的核心评分阈值与 passed 判定规则",
            font_size=16, bold=True, color=C_PRIMARY, alignment=PP_ALIGN.CENTER)

# 高级设置截图
img_path = os.path.join(SHOT_DIR, "05-advanced-settings.png")
add_image_safe(slide, img_path, Inches(2), Inches(6.1),
               Inches(9), Inches(1.2))

add_page_number(slide, 7, TOTAL_SLIDES)

# ==================== Slide 8: 技术实现 ====================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)
add_page_header(slide, "技术实现", "Tech Stack & Key Features")

# 左列
add_textbox(slide, Inches(0.6), Inches(1.7), Inches(6), Inches(0.5),
            "技术栈", font_size=20, bold=True, color=C_PRIMARY)

tech_items = [
    "Next.js 14 App Router + TypeScript",
    "Tailwind CSS 响应式样式",
    "OpenAI SDK（支持自定义 baseURL）",
    "Server-Sent Events 流式推送",
    "Electron 31 + electron-builder",
    "Next.js standalone 独立部署",
]
add_bullet_list(slide, Inches(0.6), Inches(2.3), Inches(6), Inches(3),
                tech_items, font_size=15)

# 右列
add_textbox(slide, Inches(7), Inches(1.7), Inches(6), Inches(0.5),
            "关键工程特性", font_size=20, bold=True, color=C_PRIMARY)

eng_items = [
    "词格字数最高优先（所有 prompt 第一条）",
    "纯字数修复独立 5 轮，不计入主迭代",
    "超时自动重试 5 次（指数退避 1→2→4→8→16s）",
    "AbortController 中止生成",
    "自动扩写提示词（✦ 帮写按钮）",
    "单句点击修改（保持字数不变）",
]
add_bullet_list(slide, Inches(7), Inches(2.3), Inches(6), Inches(3),
                eng_items, font_size=15)

# 支持的服务商
add_textbox(slide, Inches(0.6), Inches(5.5), Inches(12), Inches(0.4),
            "支持的大模型服务商", font_size=18, bold=True, color=C_PRIMARY)

providers = [
    "Pollinations（免费免注册）",
    "OpenAI · DeepSeek · Kimi · 通义千问",
    "智谱 GLM · SiliconFlow · OpenRouter",
    "商汤日日新 SenseNova（公测免费）",
    "自定义 OpenAI 兼容接口",
]
add_bullet_list(slide, Inches(0.6), Inches(6.0), Inches(12), Inches(1.5),
                providers, font_size=14)

add_page_number(slide, 8, TOTAL_SLIDES)

# ==================== Slide 9: 桌面安装包 ====================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)
add_page_header(slide, "桌面安装包", "Windows .exe · 开箱即用")

# 左列
add_textbox(slide, Inches(0.6), Inches(1.7), Inches(6), Inches(0.5),
            "无需预装任何依赖", font_size=20, bold=True, color=C_PRIMARY)

feature_items = [
    "✅ 双击 .exe 直接安装运行",
    "✅ 内置 Node.js 运行时",
    "✅ 内置 Chromium 浏览器引擎",
    "✅ 无 120 秒超时限制（本地服务）",
    "✅ 本地 API Key 存储，更安全",
    "✅ 桌面快捷方式自动创建",
    "✅ 支持卸载清理",
]
add_bullet_list(slide, Inches(0.6), Inches(2.3), Inches(6), Inches(3),
                feature_items, font_size=15)

# 右列 - 安装包信息
card = add_card(slide, Inches(7), Inches(1.7), Inches(5.8), Inches(3.5),
                C_CARD, C_PRIMARY)

add_textbox(slide, Inches(7.3), Inches(1.9), Inches(5), Inches(0.5),
            "📦 安装包信息", font_size=18, bold=True, color=C_PRIMARY)

pkg_info = [
    ("文件名", "小词格 Setup 1.0.0.exe"),
    ("大小", "~122 MB"),
    ("格式", "NSIS 安装程序"),
    ("架构", "x64"),
    ("版本", "v1.0.0"),
]
for i, (label, val) in enumerate(pkg_info):
    top = Inches(2.5 + i * 0.45)
    add_textbox(slide, Inches(7.3), top, Inches(1.5), Inches(0.4),
                label, font_size=14, color=C_TEXT_MUTED)
    add_textbox(slide, Inches(8.8), top, Inches(3.8), Inches(0.4),
                val, font_size=14, bold=True, color=C_TEXT)

# 哔哩哔哩主页集成
add_textbox(slide, Inches(0.6), Inches(5.3), Inches(6), Inches(0.4),
            "作者主页集成", font_size=18, bold=True, color=C_PRIMARY)

bili_items = [
    "📺 安装程序品牌栏显示作者主页",
    "📺 应用菜单：帮助 → 作者哔哩哔哩主页",
    "📺 关于对话框带「打开主页」按钮",
    "📺 卸载完成弹窗显示主页链接",
]
add_bullet_list(slide, Inches(0.6), Inches(5.8), Inches(6), Inches(1.5),
                bili_items, font_size=14)

# 底部区域截图
img_path = os.path.join(SHOT_DIR, "06-footer.png")
add_image_safe(slide, img_path, Inches(7), Inches(5.3),
               Inches(5.8), Inches(1.8))

add_page_number(slide, 9, TOTAL_SLIDES)

# ==================== Slide 10: 项目总结 ====================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)
add_page_header(slide, "项目总结", "What We've Built")

# 三个亮点
highlights = [
    ("🎯", "智能生成", "四个 AI 智能体协作\n自动迭代产出高质量歌词", C_PRIMARY),
    ("📐", "词格精准", "逐音节匹配词格\n支持延音 + 与英语音节", C_GREEN),
    ("🖥️", "桌面体验", "Windows 原生安装包\nSSE 流式 + 中止控制", C_ORANGE),
]

for i, (emoji, title, desc, color) in enumerate(highlights):
    left = Inches(0.6 + i * 4.2)
    top = Inches(2.0)
    w = Inches(4.0)
    h = Inches(2.5)

    card = add_card(slide, left, top, w, h, C_CARD, color)

    add_textbox(slide, left + Inches(0.3), top + Inches(0.3),
                Inches(1), Inches(1),
                emoji, font_size=40, color=color)
    add_textbox(slide, left + Inches(1.2), top + Inches(0.5),
                w - Inches(1.5), Inches(0.6),
                title, font_size=22, bold=True, color=C_TEXT)
    add_textbox(slide, left + Inches(0.3), top + Inches(1.2),
                w - Inches(0.6), Inches(1.2),
                desc, font_size=14, color=C_TEXT_DIM)

# 数据亮点
stats = [
    ("4", "AI 智能体"),
    ("10", "严格度等级"),
    ("6", "评审维度"),
    ("5+5", "字数修复轮次"),
    ("122MB", "安装包大小"),
]

for i, (num, label) in enumerate(stats):
    left = Inches(0.6 + i * 2.4)
    top = Inches(5.0)
    add_textbox(slide, left, top, Inches(2.2), Inches(0.7),
                num, font_size=36, bold=True, color=C_PRIMARY,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, top + Inches(0.7), Inches(2.2), Inches(0.4),
                label, font_size=13, color=C_TEXT_DIM,
                alignment=PP_ALIGN.CENTER)

add_page_number(slide, 10, TOTAL_SLIDES)

# ==================== Slide 11: 演示截图集 ====================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)
add_page_header(slide, "实机演示", "Real Product Screenshots")

# 2x3 截图网格
shots = [
    ("01-main-ui.png", "主界面"),
    ("02-advanced-editor.png", "高级编辑"),
    ("03-rhythm-tags.png", "节奏标记"),
    ("04-english-mode.png", "英语模式"),
    ("05-advanced-settings.png", "高级设置"),
    ("06-footer.png", "底部区域"),
]

for i, (fname, label) in enumerate(shots):
    col = i % 3
    row = i // 3
    left = Inches(0.6 + col * 4.15)
    top = Inches(1.8 + row * 2.7)

    img_path = os.path.join(SHOT_DIR, fname)
    add_image_safe(slide, img_path, left, top, Inches(4), Inches(2.2))
    add_textbox(slide, left, top + Inches(2.3), Inches(4), Inches(0.35),
                label, font_size=13, bold=True, color=C_TEXT,
                alignment=PP_ALIGN.CENTER)

add_page_number(slide, 11, TOTAL_SLIDES)

# ==================== Slide 12: 致谢 ====================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)

# 大色块装饰
deco = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
    Inches(13.333), Inches(0.15)
)
deco.fill.solid()
deco.fill.fore_color.rgb = C_PRIMARY
deco.line.fill.background()

deco2 = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.35),
    Inches(13.333), Inches(0.15)
)
deco2.fill.solid()
deco2.fill.fore_color.rgb = C_PRIMARY
deco2.line.fill.background()

add_textbox(slide, Inches(0), Inches(2.0), Inches(13.333), Inches(1),
            "感谢观看", font_size=56, bold=True, color=C_TEXT,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.6),
            "Thank You for Your Time", font_size=24, color=C_PRIMARY,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0), Inches(4.2), Inches(13.333), Inches(0.5),
            "欢迎关注我的哔哩哔哩频道",
            font_size=18, color=C_TEXT_DIM,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0), Inches(4.8), Inches(13.333), Inches(0.6),
            "space.bilibili.com/527121484",
            font_size=28, bold=True, color=C_PRIMARY,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0), Inches(6.2), Inches(13.333), Inches(0.4),
            "小词格 · AI 歌词生成器  |  v1.0.0  |  2026",
            font_size=14, color=C_TEXT_MUTED,
            alignment=PP_ALIGN.CENTER)

# ==================== 保存 ====================

os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)
prs.save(OUT_PATH)
print(f"✅ PPT 已生成: {OUT_PATH}")
print(f"   共 {len(prs.slides)} 页")
